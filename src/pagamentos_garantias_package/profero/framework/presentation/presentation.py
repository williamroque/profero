"""Esse módulo contém a classe `Presentation`."""

import pptx


class Presentation():
    """
    Essa classe controla todas as instâncias `Slide` e representa a apresentação
    inteira.
    """
    def __init__(self, inputs, props):
        self.inputs = inputs

        # Armazena todos as instâncias `Slide`
        self.slides = []

        # Inicializar uma instância da classe `pptx.Presentation` do `python-pptx`
        self.presentation = pptx.Presentation()
        self.presentation.slide_width = props['width']
        self.presentation.slide_height = props['height']

        # O caminho do arquivo de saída
        self.output_path = props['output_path']

    def add_slide(self, slide):
        """
        Adicionar uma instância `Slide`.

        * `slide (Slide)` -- A instância a ser adicionada.
        """
        self.slides.append(slide)

    def query(self, slide_id):
        """
        Encontrar uma instância `Slide` que tenha o id especificado.

        * `slide_id (str)` -- Código de identificação do slide.
        """

        # Usar um gerador para parar no primeiro valor válido
        search_generator = (slide for slide in self.slides if slide.slide_id == slide_id)
        try:
            return next(search_generator)
        except StopIteration:
            return None

    def render(self):
        """
        Chamar o método `render` em todas as instâncias `Slide` e salvar a
        apresentação no arquivo de saída.
        """

        for slide in self.slides:
            slide.render()

        self.presentation.save(self.output_path)
