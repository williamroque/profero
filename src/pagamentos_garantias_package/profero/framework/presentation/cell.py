"""Esse módulo contém a classe `Cell`."""

from pptx.util import Cm, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class Cell():
    """
    Essa classe é o átomo da apresentação; representa uma célula.

    Possui métodos para criar instâncias `Shape`. No PowerPoint,
    o shape é basicamente a superclasse da qual todos os
    elementos são derivados, inclusive texto e tabelas.
    """

    def __init__(self, inputs, props, cell_id, index, parent_row):
        self.inputs = inputs

        # A largura da célula
        self.width = props['width']

        # A posição horizontal da célula
        self.x_offset = props['x_offset']

        # O id da célula para uso com `query`
        self.cell_id = cell_id

        # A posição relativa da célula
        self.index = index

        # A instância `Row` que controla essa instância
        self.parent_row = parent_row

    def create_rect(self, x, y, rect_width, rect_height, fill_color=RGBColor(0xB, 0x5D, 0x77), inherit_shadow=False, show_border=False):
        """
        Criar um `Shape` retangular.

        * `x (float)` -- Posição horizontal.
        * `y (float)` -- Posição vertical.
        * `rect_width (float)` -- A largura do retângulo.
        * `rect_height (float)` -- A altura do retângulo.
        * `fill_color (RGBColor)` --  A cor de fundo.
        * `inherit_shadow (bool)` -- Herança da propriedade de sombra do `Shape` pai
        (em geral, `True` significa que tem sombra e `False` que não tem).
        * `show_border (bool)` -- Borda do retângulo.
        """

        # Criar `Shape` usando a instância `pptx.Slide` criada pela instância `Slide`
        # que controla a instância `Row` que controla essa instância
        shape = self.parent_row.parent_slide.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y,
            rect_width, rect_height
        )

        # Definir herdabilidade da sombra de acordo com o argumento acima
        shape.shadow.inherit = inherit_shadow

        # Usar fundo especificado pelo argumento acima
        self.set_fill_color(shape, fill_color)

        # Mostrar borda de acordo com o argumento acima
        if not show_border:
            shape.line.fill.background()

        return shape

    def set_fill_color(self, shape, color):
        """
        Usar fundo.

        * `shape (Shape)` -- `Shape` a ser modificado.
        * `color (RGBColor)` -- A cor a ser usada.
        """

        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    def set_text(self, shape, text, alignment=PP_ALIGN.LEFT, font_family='Calibri', font_size=Pt(18), bold=False, italic=None, color=RGBColor(0xFF, 0xFF, 0xFF), slide_link=None, vertical_anchor=MSO_ANCHOR.MIDDLE, margin_left=Cm(.25), margin_top=Cm(.25), margin_right=Cm(.25), margin_bottom=Cm(.25)):
        """
        Definir texto de uma instância `Shape`.

        * `shape (Shape)` -- O shape a ser modificado.
        * `text (str)` -- O conteúdo a ser adicionado.
        * `alignment (PP_ALIGN)` -- O alinhamento horizontal do texto.
        * `font_family (str)` -- O nome da fonte.
        * `font_size (Pt)` -- O tamanho da fonte.
        * `bold (bool)` -- Negrito.
        * `italic (bool)` -- Itálico.
        * `color (RGBColor)` -- Cor.
        * `slide_link (Slide)` -- Instância `Slide` para criar um link
        PowerPoint apontando para ele.
        * `vertical_anchor (MSO_ANCHOR)` -- Alinhamento vertical do texto.
        * `margin_left (Cm/Inches)` -- Margem esquerda.
        * `margin_top (Cm/Inches)` -- Margem de cima.
        * `margin_right (Cm/Inches)` -- Margem direita.
        * `margin_bottom (Cm/Inches)` -- Margem de baixo.
        """

        # Criar/esvaziar texto do `Shape`
        text_frame = shape.text_frame
        text_frame.clear()

        # Definir alinhamento vertical do texto
        text_frame.vertical_anchor = vertical_anchor

        # Definir margens
        text_frame.margin_left = margin_left
        text_frame.margin_top = margin_top
        text_frame.margin_right = margin_right
        text_frame.margin_bottom = margin_bottom

        # Obter primeiro paragrafo
        p = text_frame.paragraphs[0]

        # Definir alinhamento horizontal
        p.alignment = alignment

        # Definir texto
        run = p.add_run()
        run.text = text

        if slide_link != None:
            # Altura de linha
            line_height = font_size + Pt(2.5) # usar constante devido à falta de um método API

            # Criar um retângulo invisível em cima do texto (para não ter que usar o formato
            # padrão de links do PowerPoint)
            link_rect = self.create_rect(
                shape.left,
                shape.top + run.text.count('\n') * line_height + margin_top,
                shape.width,
                line_height
            )
            self.set_shape_transparency(link_rect, 100)

            # Definir a ação do link como abrir o slide alvo
            link_rect.click_action.target_slide = slide_link.slide

        font = run.font
        font.name = font_family
        font.size = font_size
        font.bold = bold
        font.italic = italic
        font.color.rgb = color
        font.underline = False

    def sub_element(self, parent, tagname, **kwargs):
        """
        Criar elemento XML do PowerPoint com um atributo arbitrário.
        """

        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)

        return element

    def set_shape_transparency(self, shape, alpha):
        """
        Definir a transparência de um `Shape`.
        """

        ts = shape.fill._xPr.solidFill
        sF = ts.get_or_change_to_srgbClr()
        sE = self.sub_element(sF, 'a:alpha', val=str((100 - alpha) * 1000))

    def render(self, slide):
        """
        Método de renderização a ser implementado por subclasses.

        * `slide (Slide)`
        """
