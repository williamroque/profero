"""Esse módulo contém a classe `Slide`."""

from pptx.util import Pt


class Slide():
    """
    Essa classe controla todas as instâncias da classe `Row` e representa o slide.
    """

    def __init__(self, inputs, slide_id, slide_type, index, background, parent_presentation):
        """
        * `inputs (dict)` -- Valores de entrada.
        * `slide_id (str)` -- ID do slide.
        * `slide_type (int)` -- O tipo do slide de acordo com o PowerPoint;
        o tipo 6 é o em branco
        * `index (int)` -- Posição relativa do slide.
        * `background (str)` -- O caminho do fundo a ser usado.
        * `parent_presentation (Presentation)` -- Presentação que contém esse slide.
        """

        self.inputs = inputs
        self.slide_id = slide_id
        self.slide_type = slide_type
        self.index = index
        self.parent_presentation = parent_presentation

        # Criar um slide novo pelo `python-pptx` de acordo com o tipo de
        # layout definido acima e salvar a instância criada da classe
        # `pptx.Slide`
        self.slide = parent_presentation.presentation.slides.add_slide(
            parent_presentation.presentation.slide_layouts[self.slide_type]
        )

        self.background = background

        if self.background is not None:
            # Criar um `Shape` de imagem usando o `self.background` e ajustar
            # para cobrir todo o slide
            background_image = self.slide.shapes.add_picture(
                self.background,
                Pt(-1), Pt(-1),
                width=parent_presentation.presentation.slide_width + Pt(2),
                height=parent_presentation.presentation.slide_height + Pt(3)
            )

            # Forma não muito intuitiva de colocar a imagem no segundo plano
            self.slide.shapes._spTree.remove(background_image._element)
            self.slide.shapes._spTree.insert(2, background_image._element)

        # Armazena todas as instâncias da classe `Row`
        self.rows = []

    def add_row(self, row):
        """
        Adicionar uma instância `Row`.

        * `row (Row)` -- A instância `Row` a ser adicionada.
        """
        self.rows.append(row)

    def query(self, row_id):
        """
        Encontrar uma instância `Row` que tenha o id especificado

        * `row_id (str)` -- O ID da fileira a ser procurada.
        """

        # Usar um gerador para parar no primeiro valor válido
        search_generator = (row for row in self.rows if row.row_id == row_id)
        try:
            return next(search_generator)
        except StopIteration:
            return None

    def render(self):
        """
        Chamar o método `render` de todas as instâncias `Row`.
        """

        for row in self.rows:
            row.render(self.slide)
