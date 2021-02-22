from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell

import locale
locale.setlocale(locale.LC_ALL, 'pt_BR')

import numpy as np


NOTE = """
➢ Valores com base em {}
➢ Contrato R$ MM: refere-se ao valor contratual na hora da compra
➢ LTV Últimas Vendas: baseado em todas as vendas dos últimos 6 meses
➢ Estoque últimas vendas: é a metragem total do estoque, multiplicado pela média do valor do m² das vendas dos últimos 6 meses
➢ OBS: diferença entre o total de unidades do empreendimento com relação ao número de vendas, se dá em razão dos distratos e das unidades revendidas
""".strip()


def set_cell_border(cell, border_color="000000", border_width='12700', sides=['left', 'right', 'top', 'bottom']):
    side_map = {
        'left': 'a:lnL',
        'right': 'a:lnR',
        'top': 'a:lnT',
        'bottom': 'a:lnB'
    }

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for line in sides:
        ln = Cell.sub_element(tcPr, side_map[line], w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = Cell.sub_element(ln, 'a:solidFill')
        srgbClr = Cell.sub_element(solidFill, 'a:srgbClr', val=border_color)
        prstDash = Cell.sub_element(ln, 'a:prstDash', val='solid')
        round_ = Cell.sub_element(ln, 'a:round')
        headEnd = Cell.sub_element(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = Cell.sub_element(ln, 'a:tailEnd', type='none', w='med', len='med')


class TableCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'table', 0,
            parent_row
        )

        self.row_count = 0
        self.column_number = 17

        self.props = props

        self.table = None

    def render(self, slide):
        table_width = Cm(33)
        table_height = self.parent_row.height - Cm(1)
        y_correction = Cm(-1)
        x_correction = Cm(-0.25)

        self.headers_height = 3
        data_height = 2 * len(self.props['empreendimentos'])
        total_height = 1

        table_rows_length = self.headers_height + data_height + total_height
        table_cols_length = len(self.props['empreendimentos'][0].keys()) - 1

        self.table = slide.shapes.add_table(
            table_rows_length,
            table_cols_length,
            self.x_offset + self.width / 2 - table_width / 2 + x_correction,
            self.parent_row.y_offset +\
                self.parent_row.height / 2 -\
                table_height / 2 +\
                y_correction,
            int(table_width), int(table_height)
        ).table
        self.table.horz_banding = False
        self.table.columns[0].width = Cm(2.5)

        self.headers = {
            'Empreendimento': [
                ['Empreendimento/\nFase'],
                ['Cidade/UF'],
                ['#\nUnidades'],
                [
                    'Média m²',
                    'Média R$/m²'
                ],
            ],
            'Obras': [
                ['Evolução\ndas\nObras'],
                ['Conclusão']
            ],
            'Vendas & Estoque': [
                ['#\nVendas'],
                ['%\nVendas'],
                ['#\nEstoque'],
                ['%\nEstoque'],
                ['Estoque\n(R$ MM)'],
                ['Estoque últimas vendas\n(R$ MM)']
            ],
            'Recebíveis': [
                ['#\nRecebíveis'],
                ['Recebíveis\n(R$ MM)'],
                ['Recebíveis\n+ Estoque\n(R$ MM)']
            ],
            'Valor imóveis & LTV': [
                ['Contrato\n(R$ MM)'],
                ['LTV\nÚltimas vendas']
            ]
        }

        self.add_headers()
        self.add_data()

    def add_headers(self):
        horizontal_offset = 0
        for header, subheaders in self.headers.items():
            header_cell = self.table.cell(0, horizontal_offset)
            header_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            set_cell_border(header_cell, 'FFFFFF', '15000')

            header_end_cell = self.table.cell(0, horizontal_offset + len(subheaders) - 1)
            header_cell.merge(header_end_cell)

            self.set_text(
                header_cell,
                header.upper(),
                alignment=PP_ALIGN.CENTER,
                font_family='Calibri',
                font_size=Pt(8),
                color=RGBColor(255, 255, 255),
                bold=True
            )
            self.set_fill_color(header_cell, RGBColor(0x00, 0x55, 0x7B))

            self.table.rows[0].height = Cm(.8)

            self.table.rows[1].height = Cm(1.5)
            self.table.rows[2].height = Cm(1.5)

            for col, subheader in enumerate(subheaders):
                for row, infraheader in enumerate(subheader[:2]):
                    infraheader_cell = self.table.cell(row + 1, horizontal_offset + col)
                    infraheader_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                    infraheader_cell.margin_left = 0
                    infraheader_cell.margin_right = 0

                    if row == 0:
                        set_cell_border(infraheader_cell, 'FFFFFF', '15000', ['top', 'left', 'right'])
                        set_cell_border(infraheader_cell, 'FFFFFF', '5000', ['bottom'])
                    elif row == 1:
                        set_cell_border(infraheader_cell, 'FFFFFF', '15000', ['bottom', 'left', 'right'])
                        set_cell_border(infraheader_cell, 'FFFFFF', '5000', ['top'])

                    self.set_text(
                        infraheader_cell,
                        infraheader,
                        alignment=PP_ALIGN.CENTER,
                        font_family='Calibri',
                        font_size=Pt(8),
                        color=RGBColor(255, 255, 255)
                    )
                    self.set_fill_color(infraheader_cell, RGBColor(0x00, 0x55, 0x7B))

                subheader_end_cell = self.table.cell(2, horizontal_offset + col)
                if len(subheader) < 2:
                    infraheader_cell.merge(subheader_end_cell)

            horizontal_offset += len(subheaders)

    def add_data(self):
        vertical_offset = self.headers_height

        format_currency = lambda x: '{:.2f}'.format(x / 1e6).replace('.', ',')
        format_2_perc = lambda x: '{:.2f}%'.format(round(x * 100, 2)).replace('.', ',')
        format_0_perc = lambda x: '{}%'.format(round(x * 100)).replace('.', ',')

        for empr in self.props['empreendimentos']:
            for i in range(self.column_number):
                self.set_fill_color(self.table.cell(vertical_offset, i), RGBColor(0xFF, 0xFF, 0xFF))
                self.set_fill_color(self.table.cell(vertical_offset + 1, i), RGBColor(0xFF, 0xFF, 0xFF))

            nome_cell = self.table.cell(vertical_offset, 0)
            self.add_data_text(nome_cell, empr['nome'])
            nome_cell.merge(
                self.table.cell(vertical_offset + 1, 0)
            )

            cidade_cell = self.table.cell(vertical_offset, 1)
            self.add_data_text(cidade_cell, empr['cidade'])
            cidade_cell.merge(
                self.table.cell(vertical_offset + 1, 1)
            )

            num_unidades_cell = self.table.cell(vertical_offset, 2)
            self.add_data_text(num_unidades_cell, empr['num-unidades'])
            num_unidades_cell.merge(
                self.table.cell(vertical_offset + 1, 2)
            )

            media_m2_cell = self.table.cell(vertical_offset, 3)
            self.add_data_text(media_m2_cell, empr['media-m2'])

            media_rs_m2_cell = self.table.cell(vertical_offset + 1, 3)
            self.add_data_text(media_rs_m2_cell, empr['media-rs-m2'])

            evolucao_cell = self.table.cell(vertical_offset, 4)
            self.add_data_text(evolucao_cell, empr['evolucao'], f=format_0_perc)
            evolucao_cell.merge(
                self.table.cell(vertical_offset + 1, 4)
            )

            conclusao_cell = self.table.cell(vertical_offset, 5)
            self.add_data_text(conclusao_cell, empr['conclusao'])
            conclusao_cell.merge(
                self.table.cell(vertical_offset + 1, 5)
            )

            num_vendas_cell = self.table.cell(vertical_offset, 6)
            self.add_data_text(num_vendas_cell, empr['num-vendas'])
            num_vendas_cell.merge(
                self.table.cell(vertical_offset + 1, 6)
            )

            perc_vendas_cell = self.table.cell(vertical_offset, 7)
            self.add_data_text(perc_vendas_cell, empr['perc-vendas'], f=format_2_perc)
            perc_vendas_cell.merge(
                self.table.cell(vertical_offset + 1, 7)
            )

            num_estoque_cell = self.table.cell(vertical_offset, 8)
            self.add_data_text(num_estoque_cell, empr['num-estoque'])
            num_estoque_cell.merge(
                self.table.cell(vertical_offset + 1, 8)
            )

            perc_estoque_cell = self.table.cell(vertical_offset, 9)
            self.add_data_text(perc_estoque_cell, empr['perc-estoque'], f=format_2_perc)
            perc_estoque_cell.merge(
                self.table.cell(vertical_offset + 1, 9)
            )

            estoque_cell = self.table.cell(vertical_offset, 10)
            self.add_data_text(estoque_cell, empr['estoque'], f=format_currency)
            estoque_cell.merge(
                self.table.cell(vertical_offset + 1, 10)
            )

            estoque_ultimas_cell = self.table.cell(vertical_offset, 11)
            self.add_data_text(estoque_ultimas_cell, empr['estoque-ultimas'], f=format_currency)
            estoque_ultimas_cell.merge(
                self.table.cell(vertical_offset + 1, 11)
            )

            num_recebiveis_cell = self.table.cell(vertical_offset, 12)
            self.add_data_text(num_recebiveis_cell, empr['num-recebiveis'])
            num_recebiveis_cell.merge(
                self.table.cell(vertical_offset + 1, 12)
            )

            recebiveis_cell = self.table.cell(vertical_offset, 13)
            self.add_data_text(recebiveis_cell, empr['recebiveis'], f=format_currency)
            recebiveis_cell.merge(
                self.table.cell(vertical_offset + 1, 13)
            )

            recebiveis_estoque_cell = self.table.cell(vertical_offset, 14)
            self.add_data_text(recebiveis_estoque_cell, empr['recebiveis-estoque'], f=format_currency)
            recebiveis_estoque_cell.merge(
                self.table.cell(vertical_offset + 1, 14)
            )

            contrato_cell = self.table.cell(vertical_offset, 15)
            self.add_data_text(contrato_cell, empr['contrato'], f=format_currency)
            contrato_cell.merge(
                self.table.cell(vertical_offset + 1, 15)
            )

            ltv_cell = self.table.cell(vertical_offset, 16)
            self.add_data_text(ltv_cell, empr['ltv'], f=format_0_perc)
            ltv_cell.merge(
                self.table.cell(vertical_offset + 1, 16)
            )

            vertical_offset += 2

        total_cell = self.table.cell(vertical_offset, 0)
        self.add_data_text(total_cell, 'Total/Média', bold=True)

        for i in range(self.column_number):
            cell = self.table.cell(vertical_offset, i)

            self.set_fill_color(cell, RGBColor(0xB8, 0xD7, 0xF0))
            set_cell_border(cell, 'FFFFFF', '0')

    def add_data_text(self, cell, value, f=str, bold=False):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self.set_text(
            cell,
            f(value),
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(8),
            color=RGBColor(0x00, 0x39, 0x60),
            bold=bold
        )
            

class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'ativos-imobiliarios', 6,
            index,
            None,
            parent_presentation
        )

        self.title = 'Ativos Imobiliários'

        self.props = props

        self.table_of_contents_slide = table_of_contents_slide

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height - Cm(.5),
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        table_row = Row(
            inputs,
            {
                'height': .75 * slide_height - note_height + Cm(.5),
                'y_offset': header_row.y_offset + header_row.height
            },
            'table', 1,
            self
        )

        table_cell = TableCell(inputs, slide_width, props, table_row)
        table_row.add_cell(table_cell)

        self.add_row(table_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': table_row.y_offset + table_row.height
            },
            'note', 2,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(
                props['date']
            ),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
