from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
import matplotlib.path
import matplotlib.patches as patches

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell

import re
import io


NOTE = """
➢ Valores com base em {}
➢ O Gatilho de Sobregarantia é calculado a partir da razão entre o saldo dos Direitos Creditórios Adimplidos e o saldo devedor dos CRI
➢ Direitos Creditórios Inadimplidos são os recebíveis cujas prestações não tenham sido pagas a partir do 91º dia a contar do respectivo vencimento
➢ Limíte de Garantia Mínima: {:.0%}
""".strip()

INFO_TEXT = """
• Sobregarantia Recebíveis: R$ {:.2f} MM / R$ {:.2f} MM = {:.2f}%
• Sobregarantia Total: R$ {:.2f} MM / R$ {:.2f} MM = {:.2f}%
"""


class ChartCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'chart', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

    def render(self, slide):
        slide = self.parent_row.parent_slide

        values = np.array([
            self.props['recebiveis-adimplidos'],
            self.props['recebiveis-inadimplidos'],
            self.props['estoque'],
            self.props['fundo-reserva'],
        ])

        dual_values = np.array([
            self.props['x-senior'],
            self.props['x-subordinada'],
        ])

        labels = (
            'RECEBÍVEIS\nADIMPLENTES',
            'RECEBÍVEIS\nINADIMPLENTES',
            'ESTOQUE',
            'FUNDO DE RESERVA',
        )

        dual_labels = [
            'SÊNIOR',
            'SUBORDINADA',
        ]

        chart_width = 10.5
        chart_height = 4

        plot_width = .7
        plot_x = 1/2 - plot_width/2

        fig = plt.figure(figsize=(chart_width, chart_height))
        ax = fig.add_axes([plot_x, .05, plot_width, .9])

        plt.axis('off')

        ax.set_xlim(-.5, .5)

        bar_width = .26
        bar_margin = .09

        bar_displacement = bar_width / 2 + bar_margin

        annotation_bracket_height = 14

        colors = [
            '#D30000',
            '#2E4F99',
        ]

        font_colors = [
            '#FFFFFF',
            '#FFFFFF',
        ]

        for (i, value), csum in zip(enumerate(dual_values), dual_values.cumsum()):
            bottom = csum - value

            ax.bar(
                -bar_displacement, value,
                bar_width, bottom=bottom,
                color=colors[i]
            )
            plt.text(
                -bar_displacement, bottom + value / 2,
                '{}\nR$ {:.2f} MM'.format(
                    dual_labels[i],
                    value / 1e+6
                ).replace('.', ','),
                ha='center',
                va='center',
                color=font_colors[i],
                fontsize=9,
                fontweight='extra bold',
                fontname='Calibri'
            )

        ax.annotate(
            'R$ {:.2f} MM'.format(
                dual_values.sum() / 1e+6
            ).replace('.', ','),
            xy=(-bar_margin + .015, dual_values.sum() / 2),
            xytext=(-bar_margin + .035, dual_values.sum() / 2),
            fontsize=9,
            fontweight='extra bold',
            fontname='Calibri',
            ha='left',
            va='center',
            arrowprops=dict(
                arrowstyle='-[, widthB={}, lengthB=.5'.format(
                    annotation_bracket_height * dual_values.sum() / values.sum()
                ),
                lw=.5
            )
        )

        colors = [
            '#00BA54',
            '#FF0000',
            '#436AC7',
            '#F8FF00',
        ]
        font_colors = [
            '#ffffff',
            '#ffffff',
            '#ffffff',
            '#003960',
        ]

        max_value = values.max()

        for (i, value), csum in zip(enumerate(values), values.cumsum()):
            bottom = csum - value

            height_threshold = .03 * max_value

            ax.bar(
                bar_displacement, value,
                bar_width, bottom=bottom,
                color=colors[i]
            )
            plt.text(
                bar_displacement, bottom + value / 2,
                '{}\nR$ {:.2f} MM'.format(
                    labels[i],
                    value / 1e+6
                ).replace('.', ','),
                ha='center',
                va='center' if value > height_threshold else 'bottom',
                color=font_colors[i],
                fontsize=9,
                fontweight='extra bold',
                fontname='Calibri'
            )

        ax.annotate(
            'R$ {:.2f} MM'.format(
                values.sum() / 1e+6
            ).replace('.', ','),
            xy=(.37, values.sum() / 2),
            xytext=(.4, values.sum() / 2),
            fontsize=9,
            fontname='Calibri',
            fontweight='extra bold',
            ha='left',
            va='center',
            arrowprops=dict(
                arrowstyle='-[, widthB={}, lengthB=.5'.format(
                    annotation_bracket_height
                ),
                lw=.5
            )
        )

        image_stream = io.BytesIO()
        fig.savefig(image_stream, format='png', dpi=300)

        chart_width = Inches(chart_width)
        chart_height = Inches(chart_height)

        slide.slide.shapes.add_picture(
            image_stream,
            self.slide_width / 2 - chart_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - chart_height / 2,
            chart_width,
            chart_height
        )


class InfoCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'info', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

    def render(self, slide):
        box_width = self.width * .5
        box_height = self.parent_row.height

        info_box = self.create_rect(
            self.x_offset + self.width / 2 - box_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - box_height,
            box_width, box_height
        )
        self.set_shape_transparency(info_box, 100)

        values = np.array([
            self.props['recebiveis-adimplidos'],
            self.props['recebiveis-inadimplidos'],
            self.props['estoque'],
            self.props['fundo-reserva'],
        ])

        dual_values = np.array([
            self.props['x-senior'],
            self.props['x-subordinada'],
        ])

        recebiveis_adimplidos = self.props['recebiveis-adimplidos']

        dual_values_sum = dual_values.sum()
        values_sum = values.sum()

        sobregarantia_recebiveis = recebiveis_adimplidos / dual_values_sum
        sobregarantia_total = values_sum / dual_values_sum

        self.set_text(
            info_box,
            INFO_TEXT.format(
                recebiveis_adimplidos / 1e+6,
                dual_values_sum / 1e+6,
                sobregarantia_recebiveis * 100,
                values_sum / 1e+6,
                dual_values_sum / 1e+6,
                sobregarantia_total * 100
            ).replace('.', ','),
            font_family='Calibri',
            font_size=Pt(14),
            color=RGBColor(0x10, 0x20, 0x30),
            alignment=PP_ALIGN.LEFT,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )
            


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'garantia', 6,
            index,
            None,
            parent_presentation,
            'Garantia',
            table_of_contents_slide
        )

        self.props = props

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        chart_row = Row(
            inputs,
            {
                'height': .7 * slide_height - note_height,
                'y_offset': header_row.y_offset + header_row.height
            },
            'chart', 1,
            self
        )

        chart_cell = ChartCell(inputs, slide_width, self.props, chart_row)
        chart_row.add_cell(chart_cell)

        self.add_row(chart_row)

        info_row = Row(
            inputs,
            {
                'height': .05 * slide_height,
                'y_offset': chart_row.y_offset + chart_row.height
            },
            'info', 2,
            self
        )

        info_cell = InfoCell(inputs, slide_width, self.props, info_row)
        info_row.add_cell(info_cell)

        self.add_row(info_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': info_row.y_offset + info_row.height
            },
            'note', 3,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(
                props['date'],
                props['garantia-minima'] / inputs.get('saldo-cri')
            ),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
