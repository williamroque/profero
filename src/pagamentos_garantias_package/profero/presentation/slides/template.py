from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell


NOTE = """
➢ <++>
""".strip()


class <++>(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            <++>, 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

    def render(self, slide):
        slide = self.parent_row.parent_slide
        slide.table_of_contents_slide.add_entry(
            slide.title, [slide.index + 1], slide
        )


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            <++>, 6,
            index,
            None,
            parent_presentation,
            <++>,
            table_of_contents_slide
        )

        self.props = props

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        <++> = Row(
            inputs,
            {
                'height': .75 * slide_height - note_height,
                'y_offset': header_row.y_offset + header_row.height
            },
            <++>, 1,
            self
        )

        <++> = <++>(inputs, slide_width, self.props, <++>)
        <++>.add_cell(<++>)

        self.add_row(<++>)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': <++>.y_offset + <++>.height
            },
            'note', 2,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(
                <++>
            ),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
