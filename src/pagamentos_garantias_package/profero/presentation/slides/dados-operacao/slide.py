from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell

import re


NOTE = """
Valores com base em {}
""".strip()


class TableCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'table', 0,
            parent_row
        )

        self.row_count = 0

        self.props = props

    def render(self, slide):
        table_width = Cm(13.06)
        table_height = Cm(11.18)
        y_correction = Cm(-1)

        primeira_serie = str(self.inputs.get('primeira-serie'))
        segunda_serie = str(self.inputs.get('primeira-serie') + 1)

        saldo_primeira = self.props[primeira_serie]['saldo-devedor']
        saldo_segunda = self.props[segunda_serie]['saldo-devedor']

        self.table = slide.shapes.add_table(
            13, 3,
            self.x_offset + self.width / 2 - table_width / 2,
            self.parent_row.y_offset +\
                self.parent_row.height / 2 -\
                table_height / 2 +\
                y_correction,
            int(table_width), int(table_height)
        ).table

        header_cell = self.table.cell(0, 0)
        self.set_text(
            header_cell,
            'Dados',
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        self.set_fill_color(header_cell, RGBColor(0x16, 0x36, 0x5C))

        primeira_serie_cell = self.table.cell(0, 1)
        self.set_text(
            primeira_serie_cell,
            '{}ª Série'.format(primeira_serie),
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        self.set_fill_color(primeira_serie_cell, RGBColor(0x16, 0x36, 0x5C))

        segunda_serie_cell = self.table.cell(0, 2)
        self.set_text(
            segunda_serie_cell,
            '{}ª Série'.format(segunda_serie),
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        self.set_fill_color(segunda_serie_cell, RGBColor(0x16, 0x36, 0x5C))

        self.add_table_row(
            'IF',
            self.props[primeira_serie]['instrumento-financeiro'],
            self.props[segunda_serie]['instrumento-financeiro']
        )
        self.add_table_row(
            'ISIN',
            self.props[primeira_serie]['isin'],
            self.props[segunda_serie]['isin']
        )
        self.add_table_row('Série', primeira_serie, segunda_serie)
        self.add_table_row(
            'Cedente',
            self.props[primeira_serie]['cedente'],
            self.props[segunda_serie]['cedente']
        )
        self.add_table_row(
            'Correção',
            self.props[primeira_serie]['correcao'],
            self.props[segunda_serie]['correcao']
        )
        self.add_table_row(
            'Juros',
            '{}%'.format(
                self.props[primeira_serie]['juros'] * 100
            ).replace('.', ','),
            '{}%'.format(
                self.props[segunda_serie]['juros'] * 100
            ).replace('.', ',')
        )
        self.add_table_row(
            'Data de Emissão',
            self.props[primeira_serie]['data-emissao'],
            self.props[segunda_serie]['data-emissao']
        )
        self.add_table_row(
            'Vencimento',
            self.props[primeira_serie]['vencimento'],
            self.props[segunda_serie]['vencimento']
        )
        self.add_table_row('Subordinação', 'Sênior', 'Subordinada')
        self.add_table_row(
            'Valor de Emissão',
            'R$ {:.2f} MM'.format(
                self.props[primeira_serie]['valor-emissao'] / 1e+6
            ).replace('.', ','),
            'R$ {:.2f} MM'.format(
                self.props[segunda_serie]['valor-emissao'] / 1e+6
            ).replace('.', ',')
        )
        self.add_table_row(
            'Saldo Devedor do CRI',
            'R$ {:.2f} MM'.format(
                saldo_primeira / 1e+6
            ).replace('.', ','),
            'R$ {:.2f} MM'.format(
                saldo_segunda / 1e+6
            ).replace('.', ',')
        )
        self.add_table_row(
            'Saldo dos CRI',
            'R$ {:.2f} MM'.format(
                self.inputs.get('saldo-cri') / 1e+6
            ).replace('.', ','),
            merge=True
        )

        slide = self.parent_row.parent_slide

        slide.table_of_contents_slide.add_entry(
            slide.title, [slide.index + 1], self.parent_row.parent_slide
        )

    def add_table_row(self, header, value_1, value_2=None, merge=False):
        header_cell = self.table.cell(self.row_count + 1, 0)
        self.set_text(
            header_cell,
            str(header),
            alignment=PP_ALIGN.LEFT,
            font_family='Calibri',
            font_size=Pt(12),
            color=RGBColor(0, 0, 0)
        )

        value_1_cell = self.table.cell(self.row_count + 1, 1)
        self.set_text(
            value_1_cell,
            str(value_1),
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            color=RGBColor(0, 0, 0)
        )

        value_2_cell = self.table.cell(self.row_count + 1, 2)
        if value_2 != None:
            self.set_text(
                value_2_cell,
                str(value_2),
                alignment=PP_ALIGN.CENTER,
                font_family='Calibri',
                font_size=Pt(12),
                color=RGBColor(0, 0, 0)
            )

        if merge:
            value_1_cell.merge(value_2_cell)

            header_cell.text_frame.paragraphs[0].runs[0].font.bold = True
            value_1_cell.text_frame.paragraphs[0].runs[0].font.bold = True

        self.row_count += 1


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'dados-operacao', 6,
            index,
            None,
            parent_presentation
        )

        self.title = 'Características da operação'

        self.table_of_contents_slide = table_of_contents_slide

        primeira_serie = str(self.inputs.get('primeira-serie'))
        segunda_serie = str(self.inputs.get('primeira-serie') + 1)

        saldo_primeira = props[primeira_serie]['saldo-devedor']
        saldo_segunda = props[segunda_serie]['saldo-devedor']
        self.inputs.update('saldo-cri', saldo_primeira + saldo_segunda)

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        table_row = Row(
            inputs,
            {
                'height': .75 * slide_height - note_height,
                'y_offset': header_row.y_offset + header_row.height
            },
            'table', 1,
            self
        )

        table_cell = TableCell(inputs, slide_width, props, table_row)

        table_row.add_cell(table_cell)

        self.add_row(table_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': table_row.y_offset + table_row.height
            },
            'note', 2,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(inputs.get('date')),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
