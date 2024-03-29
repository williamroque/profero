from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from profero.framework.presentation.cell import Cell
from profero.framework.presentation.row import Row

import importlib.resources
import profero.assets

import re


class ClientLogoCell(Cell):
    def __init__(self, inputs, client_logo_path, slide_width, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width * .15,
                'x_offset': 0
            },
            'client-logo', 0,
            parent_row
        )

        self.client_logo_path = client_logo_path

        self.picture_width = Cm(3.3)
        self.picture_height = Cm(2.55)

    def render(self, slide):
        slide.shapes.add_picture(
            self.client_logo_path,
            self.width / 2 - self.picture_width / 2,
            self.parent_row.height / 2 - self.picture_height / 2,
            self.picture_width, self.picture_height
        )


# Célula principal
class HeaderCell(Cell):
    def __init__(self, inputs, slide_width, x_offset, title, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width * .68,
                'x_offset': x_offset
            },
            'header', 1,
            parent_row
        )

        self.title = title
        self.slide_width = slide_width

    def render(self, slide):
        # Linha horizontal do cabeçalho
        line = self.create_rect(
            self.x_offset, self.parent_row.y_offset + self.parent_row.height / 2,
            self.width, Pt(2),
            RGBColor(0x00, 0x6B, 0xA2)
        )

        # Altura dos elementos centrais
        rect_height = Cm(.7)

        # Largura do retângulo do número de slide
        slide_number_width = Cm(1.1)

        # Largura do retângulo do aviso 'CONFIDENCIAL'
        confidencial_width = Cm(4)

        # Espaço entre o número e o aviso
        spacing = Cm(.4)

        # Espaço horizontal total ocupado pelos elementos + o espaço entre eles
        total_rect_width = slide_number_width + spacing + confidencial_width

        # Espaço vertical entre os elementos e a linha
        vertical_margin = Cm(.5)

        # Retângulo do número de slide
        slide_number = self.create_rect(
            self.slide_width / 2 - total_rect_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - rect_height - vertical_margin,
            slide_number_width, rect_height,
            RGBColor(0x0A, 0x56, 0x79)
        )

        self.set_text(
            slide_number,
            str(self.parent_row.parent_slide.index + 1),
            PP_ALIGN.CENTER,
            'Helvetica', Pt(10)
        )

        # Retângulo do aviso
        confidencial = self.create_rect(
            self.slide_width / 2 - total_rect_width / 2 + slide_number_width + spacing,
            self.parent_row.y_offset + self.parent_row.height / 2 - rect_height - vertical_margin,
            confidencial_width, rect_height,
            RGBColor(0x0A, 0x56, 0x79)
        )

        self.set_text(
            confidencial,
            'CONFIDENCIAL',
            PP_ALIGN.CENTER,
            'Helvetica', Pt(10)
        )

        # Dimensões da caixa de texto do título
        title_width = Cm(11.91)
        title_height = Cm(1.11)

        # Caixa de texto do título
        title = self.create_rect(
            self.slide_width / 2 - title_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 + vertical_margin,
            title_width, title_height
        )

        self.set_shape_transparency(title, 100)

        self.set_text(
            title,
            re.sub(r'_(\w+)_', lambda m: m.group(1).lower(), self.title.upper()),
            PP_ALIGN.CENTER,
            'Helvetica', Pt(16),
            True, None,
            RGBColor(0x09, 0x51, 0x72)
        )


class LogoCell(Cell):
    def __init__(self, inputs, slide_width, x_offset, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width * .17,
                'x_offset': x_offset
            },
            'logo', 2,
            parent_row
        )

        self.picture_width = Cm(4.68)
        self.picture_height = Cm(2.24)

    def render(self, slide):
        with importlib.resources.path(profero.assets, 'logo.png') as p:
            logo_path = str(p)
        slide.shapes.add_picture(
            logo_path,
            self.x_offset + self.width / 2 - self.picture_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - self.picture_height / 2,
            self.picture_width, self.picture_height
        )


# Essa classe representa o cabeçalho do slide
class HeaderRow(Row):
    def __init__(self, inputs, props, index, title, slide_width, slide_height, parent_slide):
        super().__init__(
            inputs,
            props,
            'header', index,
            parent_slide
        )

        client_logo_cell = ClientLogoCell(inputs, inputs.get('client-logo'), slide_width, self)
        self.add_cell(client_logo_cell)

        header_cell = HeaderCell(inputs, slide_width, client_logo_cell.x_offset + client_logo_cell.width, title, self)
        self.add_cell(header_cell)

        logo_cell = LogoCell(inputs, slide_width, header_cell.x_offset + header_cell.width, self)
        self.add_cell(logo_cell)
