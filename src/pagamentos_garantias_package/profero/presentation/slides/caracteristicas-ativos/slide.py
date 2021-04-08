from pptx.util import Cm, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import matplotlib.pyplot as plt
from cycler import cycler

import numpy as np

import io

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell


NOTE = """
➢ Valores com base em {}
""".strip()


COLORSCHEMES = [
    ['#CF918F', '#A74442', '#70568D', '#D7813C'],
    ['#4472C5', '#70568D', '#87A34C', '#A74442'],
    ['#91A7CD', '#4470A5', '#4095AD', '#87A34C'],
    ['#BD4E4C', '#4E80BB', '#87A34C', '#91A7CD']
]


class ChartCell(Cell):
    def __init__(self, inputs, slide_width, props, index, header, content, parent_row, width, x_offset, general_index):
        super().__init__(
            inputs,
            {
                'width': width,
                'x_offset': x_offset
            },
            'chart', index,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

        self.width = width
        self.x_offset = x_offset

        self.general_index = general_index

        self.header = header
        self.content = content

    def render(self, slide):
        labels, values = list(zip(*self.content))

        chart_width = self.width / Inches(1)
        chart_height = self.parent_row.height / Inches(1)

        plot_size = 1
        plot_x = 1/2 - plot_size/2 + .2

        fig = plt.figure(figsize=(chart_width, chart_height))
        ax = fig.add_axes([plot_x, 0, plot_size, plot_size])

        ax.text(
            -3.2, .6,
            self.header.upper(),
            {
                'family': 'Calibri',
                'color': '#003960',
                'weight': 'bold',
                'size': 20
            }
        )

        def func(pct, _):
            return '{}%'.format(int(pct))

        wedges, texts, autotexts = ax.pie(
            values,
            textprops=dict(
                color='w',
                fontname='Calibri',
            ),
            colors=COLORSCHEMES[self.general_index],
            rotatelabels=True,
            autopct=lambda pct: func(pct, values),
            explode=[0] + [.15] * (len(labels) - 1),
            startangle=90,
            counterclock=False,
        )

        ax.legend(
            wedges,
            labels,
            bbox_to_anchor=(-.3, .5),
            frameon=False
        )

        plt.setp(autotexts, size=8, weight='bold')

        image_stream = io.BytesIO()
        fig.savefig(image_stream, format='png', dpi=300)

        chart_width = Inches(chart_width)
        chart_height = Inches(chart_height)

        slide.shapes.add_picture(
            image_stream,
            self.x_offset,
            self.parent_row.y_offset + self.parent_row.height / 2 - chart_height / 2,
            chart_width,
            chart_height
        )


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'caracteristicas-ativo', 6,
            index,
            None,
            parent_presentation,
            'Características dos Ativos',
            table_of_contents_slide
        )

        self.props = props

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .2 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        row_count = len(self.props['charts'])
        chart_height = .8 * slide_height / row_count  - note_height / row_count

        general_index = 0
        for i, row in enumerate(self.props['charts']):
            chart_row = Row(
                inputs,
                {
                    'height': chart_height,
                    'y_offset': header_row.y_offset + header_row.height + i * chart_height
                },
                'chart', 1 + i,
                self
            )

            chart_width = slide_width / len(row.keys())

            for j, (header, content) in enumerate(row.items()):
                chart_cell = ChartCell(
                    inputs,
                    slide_width,
                    self.props,
                    j,
                    header,
                    content,
                    chart_row,
                    chart_width,
                    chart_width * j,
                    general_index
                )
                chart_row.add_cell(chart_cell)

                general_index += 1

            self.add_row(chart_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': chart_row.y_offset + chart_row.height
            },
            'note', 2,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(
                props['date']
            ),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
