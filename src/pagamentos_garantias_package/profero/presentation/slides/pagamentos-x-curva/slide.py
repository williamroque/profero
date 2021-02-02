from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell

from cycler import cycler

import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
import matplotlib.font_manager as font_manager

from profero.util.parser.spreadsheet_parser import Parser

import numpy as np

import io

import locale
locale.setlocale(locale.LC_ALL, 'pt_BR')


NOTE = """
➢ Curva: pagamentos previstos na curva de amortização original dos CRI
➢ Pagamentos: valores efetivamente pagos a título de amortização e juros dos CRI seniores e subordinado
➢ Valores em reais
➢ Valores são cumulativos na tabela e médios no gráfico para trimestres, semestres e anos
""".strip()


class ChartCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'chart', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

        self.merge_residue = []

        boletim_schema = {
            'file-type': 'csv',
            'sections': {
                'root': {
                    'header-row': 7,
                    'groups': {
                        'valores': {
                            'query': 'Valor\nPago',
                            'subquery': ['Identificação', 'Total Geral:'],
                            'dtype': 'float',
                        },
                    }
                }
            }
        }

        parser = Parser(boletim_schema)

        self.recebimento = []

        for month in self.props['recebimento']:
            total_geral = 0
            for path in month:
                print(path)
                total_geral += parser.read(path)['root']['valores'][0]

            self.recebimento.append(total_geral)

    def merge(self, n):
        if n == 3:
            self.merge_residue.append('Trimestre — ' + self.months[0][-2:])

            self.months[:3] = []
            self.months.append('Trimestre')
        elif n == 6:
            self.merge_residue.append('Semestre — ' + self.months[0][-2:])

            self.months[:6] = []
            self.months.append('Semestre')
        else:
            self.merge_residue.append('Ano — ' + self.months[0][-2:])

            self.months[:12] = []
            self.months.append('Ano')

    def get_overflow(self):
        max_months = 6
        return max(0, len(self.months) - max_months)

    def get_divisor(self, overflow):
        units = [3, 6, 12]
        return 1 + ((1 / min(overflow / (u - 1) for u in units if overflow / (u - 1) >= 1)) * overflow if overflow > units[0] else units[0] - 1)

    def merge_all(self):
        while (overflow := self.get_overflow()) > 0:
            self.merge(self.get_divisor(overflow))

    def get_min_index(self):
        periods = ['Trimestre', 'Semestre', 'Ano']
        indices = []

        for string in periods:
            try:
                indices.append(self.months.index(string))
            except ValueError:
                pass

        return min(indices) if len(indices) > 0 else 0

    def render(self, slide):
        primeira_serie = str(self.inputs.get('primeira-serie'))
        segunda_serie = str(self.inputs.get('primeira-serie') + 1)

        atual_inputs = zip(
            self.props[primeira_serie],
            self.props[segunda_serie]
        )

        self.labels = [
            'Curva',
            'Amortização Seniores',
            'Juros Seniores',
            'AMEX Seniores',
            'Amortização Subordinada',
            'Juros Subordinada',
            'AMEX Subordinada',
            'Pagamento',
            'Recebimento'
        ]

        self.months = []
        curva = []
        pagamento = []
        juros_sen = []
        amort_sen = []
        amex_sen = []
        juros_sub = []
        amort_sub = []
        amex_sub = []

        for sen, sub in atual_inputs:
            self.months.append(sen[0])

            curva.append(
                sen[1] + sen[2] +\
                sub[1] + sub[2]
            )
            pagamento.append(
                sen[1] + sen[2] + sen[3] +\
                sub[1] + sub[2] + sub[3]
            )

            juros_sen.append(sen[1])
            juros_sub.append(sub[1])

            amort_sen.append(sen[2])
            amort_sub.append(sub[2])

            amex_sen.append(sen[3])
            amex_sub.append(sub[3])

        current_quarter = self.months[-3:]
        del self.months[-3:]

        self.merge_all()
        merged_index = self.get_min_index()
        self.months = self.months[merged_index:] + self.months[:merged_index]

        self.months += current_quarter

        long_term_count = 0
        for i, term in enumerate(self.months):
            if term == 'Trimestre':
                curva[i:i+3] = [sum(curva[i:i+3]) / 3]
                pagamento[i:i+3] = [sum(pagamento[i:i+3]) / 3]
                juros_sen[i:i+3] = [sum(juros_sen[i:i+3]) / 3]
                juros_sub[i:i+3] = [sum(juros_sub[i:i+3]) / 3]
                amort_sen[i:i+3] = [sum(amort_sen[i:i+3]) / 3]
                amort_sub[i:i+3] = [sum(amort_sub[i:i+3]) / 3]
                amex_sen[i:i+3] = [sum(amex_sen[i:i+3]) / 3]
                amex_sub[i:i+3] = [sum(amex_sub[i:i+3]) / 3]
                self.recebimento[i:i+3] = [sum(self.recebimento[i:i+3]) / 3]

                self.months[i] = '{}º {}'.format(
                    self.merge_residue[:long_term_count].count(
                        self.merge_residue[long_term_count]
                    ) + 1,
                    self.merge_residue[long_term_count]
                )

                long_term_count += 1
            elif term == 'Semestre':
                curva[i:i+6] = [sum(curva[i:i+6]) / 6]
                pagamento[i:i+6] = [sum(pagamento[i:i+6]) / 6]
                juros_sen[i:i+6] = [sum(juros_sen[i:i+6]) / 6]
                juros_sub[i:i+6] = [sum(juros_sub[i:i+6]) / 6]
                amort_sen[i:i+6] = [sum(amort_sen[i:i+6]) / 6]
                amort_sub[i:i+6] = [sum(amort_sub[i:i+6]) / 6]
                amex_sen[i:i+6] = [sum(amex_sen[i:i+6]) / 6]
                amex_sub[i:i+6] = [sum(amex_sub[i:i+6]) / 6]
                self.recebimento[i:i+6] = [sum(self.recebimento[i:i+6]) / 6]

                self.months[i] = '{}º {}'.format(
                    self.merge_residue[:long_term_count].count(
                        self.merge_residue[long_term_count]
                    ) + 1,
                    self.merge_residue[long_term_count]
                )

                long_term_count += 1
            elif term == 'Ano':
                curva[i:i+12] = [sum(curva[i:i+12]) / 12]
                pagamento[i:i+12] = [sum(pagamento[i:i+12]) / 12]
                juros_sen[i:i+12] = [sum(juros_sen[i:i+12]) / 12]
                juros_sub[i:i+12] = [sum(juros_sub[i:i+12]) / 12]
                amort_sen[i:i+12] = [sum(amort_sen[i:i+12]) / 12]
                amort_sub[i:i+12] = [sum(amort_sub[i:i+12]) / 12]
                amex_sen[i:i+12] = [sum(amex_sen[i:i+12]) / 12]
                amex_sub[i:i+12] = [sum(amex_sub[i:i+12]) / 12]
                self.recebimento[i:i+12] = [sum(self.recebimento[i:i+12]) / 12]

                self.months[i] = '{}º {}'.format(
                    self.merge_residue[:long_term_count].count(
                        self.merge_residue[long_term_count]
                    ) + 1,
                    self.merge_residue[long_term_count]
                )

                long_term_count += 1

        self.values = np.array(list(zip(
            curva,
            amort_sen,
            juros_sen,
            amex_sen,
            amort_sub,
            juros_sub,
            amex_sub,
            pagamento,
            self.recebimento
        )))

        chart_width = 12.5
        chart_height = self.parent_row.height / Inches(1)

        plot_width = .75
        plot_x = 1/2 - plot_width/2 + .1

        plot_height = .8
        plot_y = .1 # unintuitive, but positive values offset upwards (positive Cartesian coordinates)

        fig = plt.figure(figsize=(chart_width, chart_height))
        ax = fig.add_axes([plot_x, plot_y, plot_width, plot_height])

        colors_cycler = cycler(color=[
            '#426AC7',
            '#A5A5A5',
            '#4C98D8',
            '#243F7A',
            '#636363',
            '#145B93'
        ])
        markers_cycler = cycler(marker=['o', '^'])

        ax.set_prop_cycle(markers_cycler * colors_cycler)

        ax.plot(self.values)

        max_val = self.values.max()
        power = 10**np.floor(np.log10(max_val))
        ceiling = np.ceil(max_val / power) * power
        step = ceiling // 10

        ax.ticklabel_format(useOffset=False, style='plain') 
        ax.yaxis.set_major_formatter(
            FuncFormatter(lambda x, _: '{:,.2f}'.format(x) if x != 0 else '-')
        )

        ax.tick_params(
            axis='both',
            colors='#0F3B5E',
            labelsize=7
        )

        ax.grid(
            axis='y',
            color='#eee'
        )
        ax.set_axisbelow(True)

        ax.spines['bottom'].set_color('#ddd')
        ax.spines['top'].set_color('#ddd')
        ax.spines['left'].set_color('#666')
        ax.spines['right'].set_color('#fff')

        ax.set_yticks(np.arange(0, ceiling + 1, step))
        ax.set_xticks(np.arange(len(self.months)))
        ax.set_xticklabels(self.months)

        font = font_manager.FontProperties(
            family='helvetica',
            weight=800,
            size=8
        )

        legend = ax.legend(
            prop=font,
            labels=self.labels,
            bbox_to_anchor=(-.1, .92),
            frameon=False
        )

        for text in legend.get_texts():
            text.set_color('#0F3B5E')

        image_stream = io.BytesIO()
        fig.savefig(image_stream, format='png', dpi=300)

        chart_width = Inches(chart_width)
        chart_height = Inches(chart_height)

        slide.shapes.add_picture(
            image_stream,
            self.slide_width / 2 - chart_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - chart_height / 2,
            chart_width,
            chart_height
        )

        slide = self.parent_row.parent_slide
        slide.table_of_contents_slide.add_entry(
            slide.title, [slide.index + 1], slide
        )


class TableCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'table', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

        self.row_count = 0

    def render(self, slide):
        table_width = Cm(30)
        table_height = self.parent_row.height - Cm(1)
        y_correction = Cm(-1)

        chart_cell = self.parent_row.parent_slide.query('chart').query('chart')
        headers = chart_cell.labels
        values = chart_cell.values
        months = chart_cell.months

        self.table = slide.shapes.add_table(
            len(headers) + 1,
            len(months) + 1,
            self.x_offset + self.width / 2 - table_width / 2,
            self.parent_row.y_offset +\
                self.parent_row.height / 2 -\
                table_height / 2 +\
                y_correction,
            int(table_width), int(table_height)
        ).table

        self.add_table_row(
            ['Tipo'] + months,
            bold=True,
            color=RGBColor(255, 255, 255),
            fill_color=RGBColor(0x16, 0x36, 0x5C),
            font_size=Pt(9)
        )

        def get_multiplier(month):
            if 'Trimestre' in month:
                return 3
            elif 'Semestre' in month:
                return 6
            elif 'Ano' in month:
                return 12
            else:
                return 1

        for header, *row in np.concatenate(([headers], values)).T:
            self.add_table_row(
                (
                    (locale.format_string(
                        'R$ %.3f',
                        float(value) * get_multiplier(months[i]),
                        True
                    ) if float(value) > 0 else '-')
                    for i, value in enumerate(row)
                ), header
            )

    def add_table_row(self, values, header=None, bold=False, color=RGBColor(0x0F, 0x3B, 0x5E), fill_color=None, font_size=Pt(9)):
        cell_offset = 0

        if header != None:
            cell = self.table.cell(self.row_count, cell_offset)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            self.set_text(
                cell,
                str(header),
                alignment=PP_ALIGN.LEFT,
                font_family='Calibri',
                font_size=font_size - Pt(1),
                color=color,
                bold=True
            )

            cell_offset += 1

        for value_i, value in enumerate(values):
            cell = self.table.cell(self.row_count, cell_offset)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            self.set_text(
                cell,
                str(value),
                alignment=PP_ALIGN.CENTER,
                font_family='Calibri',
                font_size=font_size,
                color=color,
                bold=bold
            )

            if fill_color != None:
                self.set_fill_color(cell, fill_color)

            cell_offset += 1

        self.row_count += 1


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'pagamentos-x-curva', 6,
            index,
            None,
            parent_presentation
        )

        self.title = 'Pagamentos _x_ Curva'

        self.props = props

        self.table_of_contents_slide = table_of_contents_slide

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .2 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        chart_row = Row(
            inputs,
            {
                'height': .375 * slide_height - note_height * 2/3,
                'y_offset': header_row.y_offset + header_row.height
            },
            'chart', 1,
            self
        )

        chart_cell = ChartCell(inputs, slide_width, self.props, chart_row)
        chart_row.add_cell(chart_cell)

        self.add_row(chart_row)

        offset_adjustment = Cm(1.2)

        table_row = Row(
            inputs,
            {
                'height': .395 * slide_height - note_height * 1/3,
                'y_offset': chart_row.y_offset + chart_row.height + offset_adjustment
            },
            'table', 2,
            self
        )

        table_cell = TableCell(inputs, slide_width, self.props, table_row)
        table_row.add_cell(table_cell)

        self.add_row(table_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': table_row.y_offset + table_row.height - offset_adjustment / 2
            },
            'note', 3,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE,
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
