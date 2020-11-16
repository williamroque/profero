from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell

import locale
locale.setlocale(locale.LC_ALL, 'pt_BR')


NOTE = """
➢ Valores com base em {}
""".strip()


class TableCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'table', 0,
            parent_row
        )

        self.row_count = 0

        self.props = props

    def render(self, slide):
        table_width = Cm(15.06)
        table_height = Cm(11.18)
        y_correction = Cm(-1)

        primeira_serie = str(self.inputs.get('primeira-serie'))
        segunda_serie = str(self.inputs.get('primeira-serie') + 1)

        self.table = slide.shapes.add_table(
            14, 3,
            self.x_offset + self.width / 2 - table_width / 2,
            self.parent_row.y_offset +\
                self.parent_row.height / 2 -\
                table_height / 2 +\
                y_correction,
            int(table_width), int(table_height)
        ).table

        header_cell = self.table.cell(0, 0)
        self.set_text(
            header_cell,
            'Dados',
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        self.set_fill_color(header_cell, RGBColor(0x16, 0x36, 0x5C))

        primeira_serie_cell = self.table.cell(0, 1)
        self.set_text(
            primeira_serie_cell,
            '{}ª Série'.format(primeira_serie),
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        self.set_fill_color(primeira_serie_cell, RGBColor(0x16, 0x36, 0x5C))

        segunda_serie_cell = self.table.cell(0, 2)
        self.set_text(
            segunda_serie_cell,
            '{}ª Série'.format(segunda_serie),
            alignment=PP_ALIGN.CENTER,
            font_family='Calibri',
            font_size=Pt(12),
            bold=True,
            color=RGBColor(255, 255, 255)
        )
        self.set_fill_color(segunda_serie_cell, RGBColor(0x16, 0x36, 0x5C))

        quantidade_primeira = self.props[primeira_serie]['quantidade']
        quantidade_segunda = self.props[segunda_serie]['quantidade']

        amort_unit_primeira = self.props[primeira_serie]['amortizacao-unitaria']
        amort_unit_segunda = self.props[segunda_serie]['amortizacao-unitaria']

        amex_unit_primeira = self.props[primeira_serie]['amex-unitaria']
        amex_unit_segunda = self.props[segunda_serie]['amex-unitaria']

        pmt_unitario_primeira = self.props[primeira_serie]['pagamento-total-unidade']
        pmt_unitario_segunda = self.props[segunda_serie]['pagamento-total-unidade']

        self.add_table_row(
            'Quantidade de CRI integralizado',
            [
                '{:n}'.format(quantidade_primeira),
                '{:n}'.format(quantidade_segunda),
            ],
            [False],
        )
        self.add_table_row(
            'Juros Unitários',
            [
                locale.format_string(
                    'R$ %.8f',
                    self.props[primeira_serie]['juros-unitarios'],
                    True
                ),
                locale.format_string(
                    'R$ %.8f',
                    self.props[segunda_serie]['juros-unitarios'],
                    True
                ),
            ],
            [False],
        )
        self.add_table_row(
            'Amortização Unitária',
            [
                locale.format_string(
                    'R$ %.8f',
                    amort_unit_primeira,
                    True
                ) if amort_unit_primeira > 0 else '-',
                locale.format_string(
                    'R$ %.8f',
                    amort_unit_segunda,
                    True
                ) if amort_unit_segunda > 0 else '-',
            ],
            [False],
        )
        self.add_table_row(
            'Amortização Extraordinária Unitária',
            [
                locale.format_string(
                    'R$ %.8f',
                    amex_unit_primeira,
                    True
                ) if amex_unit_primeira > 0 else '-',
                locale.format_string(
                    'R$ %.8f',
                    amex_unit_segunda,
                    True
                ) if amex_unit_segunda > 0 else '-',
            ],
            [False],
        )
        self.add_table_row(
            'Pagamento Total por Unidade',
            [
                locale.format_string(
                    'R$ %.8f',
                    pmt_unitario_primeira,
                    True
                ),
                locale.format_string(
                    'R$ %.8f',
                    pmt_unitario_segunda,
                    True
                ),
            ],
            [False],
        )
        self.add_table_row(
            'Pagamento Total do CRI em Circulação',
            [
                locale.format_string(
                    'R$ %.2f',
                    pmt_unitario_primeira * quantidade_primeira,
                    True
                ),
                locale.format_string(
                    'R$ %.2f',
                    pmt_unitario_segunda * quantidade_segunda,
                    True
                ),
            ],
            [False],
        )
        self.add_table_row(
            'Saldo Devedor Unitário',
            [
                locale.format_string(
                    'R$ %.2f',
                    self.inputs.get('saldo-primeira') / quantidade_primeira,
                    True
                ),
                locale.format_string(
                    'R$ %.2f',
                    self.inputs.get('saldo-segunda') / quantidade_segunda,
                    True
                ),
            ],
            [False],
        )
        self.add_table_row(
            'Saldo Devedor Total',
            [
                locale.format_string(
                    'R$ %.2f',
                    self.inputs.get('saldo-primeira'),
                    True
                ),
                locale.format_string(
                    'R$ %.2f',
                    self.inputs.get('saldo-segunda'),
                    True
                ),
                locale.format_string(
                    'R$ %.2f',
                    self.inputs.get('saldo-cri'),
                    True
                ),
                None,
            ],
            [False, True],
        )
        self.add_table_row(
            'Pagamento aos Investidores',
            [
                locale.format_string(
                    'R$ %.2f',
                    self.props[primeira_serie]['pagamento-investidores'],
                    True
                ),
                locale.format_string(
                    'R$ %.2f',
                    self.props[segunda_serie]['pagamento-investidores'],
                    True
                ),
            ],
            [False]
        )

        investidores = zip(
            self.props[primeira_serie]['investidores'],
            self.props[segunda_serie]['investidores'],
        )
        for investidor_i, investidor in enumerate(investidores):
            self.add_table_row(
                'Investidor {:02}'.format(investidor_i + 1),
                [*map(
                    lambda n: locale.format_string('R$ %.2f', n, True) if n > 0 else '-',
                    investidor
                )],
                [False]
            )

        slide = self.parent_row.parent_slide
        slide.table_of_contents_slide.add_entry(
            slide.title, [slide.index + 1], slide
        )

    def add_table_row(self, header, values, merge_row_cells):
        header_cell = self.table.cell(self.row_count + 1, 0)
        header_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        self.set_text(
            header_cell,
            str(header),
            alignment=PP_ALIGN.LEFT,
            font_family='Calibri',
            font_size=Pt(9),
            color=RGBColor(0, 0, 0),
        )

        for i in range(len(values) // 2 - 1):
            cell = self.table.cell(self.row_count + i + 2, 0)
            header_cell.merge(cell)

            header_cell = cell

        for value_i, (value_1, value_2) in enumerate(zip(values[::2], values[1::2])):
            self.row_count += 1

            value_1_cell = self.table.cell(self.row_count, 1)
            value_1_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            self.set_text(
                value_1_cell,
                str(value_1),
                alignment=PP_ALIGN.CENTER,
                font_family='Calibri',
                font_size=Pt(10),
                color=RGBColor(0, 0, 0)
            )

            value_2_cell = self.table.cell(self.row_count, 2)
            value_2_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if value_2 != None:
                self.set_text(
                    value_2_cell,
                    str(value_2),
                    alignment=PP_ALIGN.CENTER,
                    font_family='Calibri',
                    font_size=Pt(10),
                    color=RGBColor(0, 0, 0)
                )

            if merge_row_cells[value_i]:
                value_1_cell.merge(value_2_cell)


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'pagamento-investidores', 6,
            index,
            None,
            parent_presentation
        )

        self.title = 'Pagamento aos Investidores'

        self.props = props

        self.table_of_contents_slide = table_of_contents_slide

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        table_row = Row(
            inputs,
            {
                'height': .75 * slide_height - note_height,
                'y_offset': header_row.y_offset + header_row.height
            },
            'table', 1,
            self
        )

        table_cell = TableCell(inputs, slide_width, props, table_row)
        table_row.add_cell(table_cell)

        self.add_row(table_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': table_row.y_offset + table_row.height
            },
            'note', 2,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(
                props['date']
            ),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
