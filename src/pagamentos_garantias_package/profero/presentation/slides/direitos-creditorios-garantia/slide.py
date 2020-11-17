from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow
from profero.presentation.slides.common.note import NoteCell

import matplotlib.pyplot as plt

import io
import locale
locale.setlocale(locale.LC_ALL, 'pt_BR')


NOTE = """
➢ Valores com base em {}
➢ Direitos Creditórios Inadimplidos são os recebíveis cujas prestações não tenham sido pagas a partir do 91º dia a contar do respectivo vencimento
""".strip()


class TableCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'table', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

        self.row_count = 0

    def render(self, slide):
        table_width = Cm(32.55)
        table_height = Cm(4.11)
        y_correction = Cm(-1)

        headers = [
            'Empreendimentos',
            '# Contratos',
            '# Direitos Creditórios Adimplidos',
            '# Direitos Creditórios Inadimplidos',
            'Direitos Creditórios Inadimplidos (R$ MM)',
            '# Direitos Creditórios Inadimplidos (R$ MM)',
            'Total dos Direitos Creditórios (R$ MM)',
        ]

        self.table = slide.shapes.add_table(
            len(self.props['empreendimentos']) + 2,
            len(headers),
            self.x_offset + self.width / 2 - table_width / 2,
            self.parent_row.y_offset +\
                self.parent_row.height / 2 -\
                table_height / 2 +\
                y_correction,
            int(table_width), int(table_height)
        ).table

        self.add_table_row(
            headers,
            bold=True,
            color=RGBColor(255, 255, 255),
            fill_color=RGBColor(0x16, 0x36, 0x5C),
            font_size=Pt(11)
        )

        self.table.rows[0].height = Cm(1.6)

        for empreendimento_i, empreendimento in enumerate(self.props['empreendimentos']):
            direitos_adimplidos = self.props['direitos-adimplidos'][empreendimento_i]
            direitos_inadimplidos = self.props['direitos-inadimplidos'][empreendimento_i]

            self.add_table_row([
                empreendimento,
                *map(
                    lambda n: locale.format_string('%.2f', float(n), True),
                    [
                        self.props['contratos'][empreendimento_i],
                        self.props['num-direitos-adimplidos'][empreendimento_i],
                        self.props['num-direitos-inadimplidos'][empreendimento_i],
                        direitos_adimplidos,
                        direitos_inadimplidos,
                        direitos_adimplidos + direitos_inadimplidos
                    ]
                )
            ])

        self.add_table_row(
            [
                'Total',
                *map(
                    lambda n: locale.format_string('%.2f', float(n), True),
                    [
                        sum(self.props['contratos']),
                        sum(self.props['num-direitos-adimplidos']),
                        sum(self.props['num-direitos-inadimplidos']),
                        sum(self.props['direitos-adimplidos']),
                        sum(self.props['direitos-inadimplidos']),
                        sum(self.props['direitos-adimplidos']) + sum(self.props['direitos-inadimplidos']),
                    ]
                )
            ],
            bold=True
        )

        slide = self.parent_row.parent_slide
        slide.table_of_contents_slide.add_entry(
            slide.title, [slide.index + 1], slide
        )

    def add_table_row(self, values, bold=False, color=RGBColor(0x0F, 0x3B, 0x5E), fill_color=None, font_size=Pt(12)):
        for value_i, value in enumerate(values):
            cell = self.table.cell(self.row_count, value_i)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            self.set_text(
                cell,
                str(value),
                alignment=PP_ALIGN.CENTER,
                font_family='Calibri',
                font_size=font_size,
                color=color,
                bold=bold
            )

            if fill_color != None:
                self.set_fill_color(cell, fill_color)

        self.row_count += 1


class ChartCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'chart', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

    def render(self, slide):
        values = (
            sum(self.props['direitos-adimplidos']),
            sum(self.props['direitos-inadimplidos']),
        )

        labels = (
            'Direitos Creditórios Adimplidos',
            'Direitos Creditórios Inadimplidos',
        )

        chart_width = 7.63
        chart_height = self.parent_row.height / Inches(1)

        plot_size = 1
        plot_x = 1/2 - plot_size/2 + .2

        fig = plt.figure(figsize=(chart_width, chart_height))
        ax = fig.add_axes([plot_x, -.05, plot_size, plot_size])

        def func(pct, allvals):
            absolute = (pct / 100) * sum(allvals)
            return 'R$ {:.2f} MM\n{:d}%'.format(absolute, int(pct))

        wedges, texts, autotexts = ax.pie(
            values,
            textprops=dict(color='w'),
            colors=['#333F50', '#8497B0', '#ADB9CA'],
            rotatelabels=True,
            autopct=lambda pct: func(pct, values),
            explode=[0, .15],
            startangle=90,
            counterclock=False
        )

        ax.legend(
            wedges,
            labels,
            bbox_to_anchor=(-.3, .5),
            frameon=False
        )

        plt.setp(autotexts, size=8, weight='bold')

        image_stream = io.BytesIO()
        fig.savefig(image_stream, format='png')

        chart_width = Inches(chart_width)
        chart_height = Inches(chart_height)

        slide.shapes.add_picture(
            image_stream,
            self.slide_width / 2 - chart_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - chart_height / 2,
            chart_width,
            chart_height
        )



class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'direitos-creditorios-garantia', 6,
            index,
            None,
            parent_presentation
        )

        self.title = 'Direitos Creditórios em Garantia'

        self.props = props

        self.table_of_contents_slide = table_of_contents_slide

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        note_height = Cm(2.04)

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        table_row = Row(
            inputs,
            {
                'height': .35 * slide_height - note_height / 2,
                'y_offset': header_row.y_offset + header_row.height
            },
            'table', 1,
            self
        )

        table_cell = TableCell(inputs, slide_width, self.props, table_row)
        table_row.add_cell(table_cell)

        self.add_row(table_row)

        chart_row = Row(
            inputs,
            {
                'height': .4 * slide_height - note_height / 2,
                'y_offset': table_row.y_offset + table_row.height
            },
            'chart', 2,
            self
        )

        chart_cell = ChartCell(inputs, slide_width, self.props, chart_row)
        chart_row.add_cell(chart_cell)

        self.add_row(chart_row)

        note_row = Row(
            inputs,
            {
                'height': note_height,
                'y_offset': chart_row.y_offset + chart_row.height
            },
            'note', 3,
            self
        )

        note_cell = NoteCell(
            inputs,
            slide_width,
            NOTE.format(
                props['date']
            ),
            note_row
        )
        note_row.add_cell(note_cell)

        self.add_row(note_row)
