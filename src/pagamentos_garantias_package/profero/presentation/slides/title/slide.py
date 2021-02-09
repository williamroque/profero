"""
Esse slide representa o título da apresentação
"""

import importlib.resources

import time
import locale

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell

import profero.assets

# Definir a região como Brasil para uso com métodos de formatação numérica
locale.setlocale(locale.LC_TIME, 'pt_BR')

# Modelo do título
TITLE_STRING = """
CRI Logos {}ª e {}ª Séries – Relatório de Pagamentos e Garantias
Certificados de Recebíveis Imobiliários
Setembro de 2020
""".strip()


class ClientLogoCell(Cell):
    """
    Célula responsável pelo logotipo do cliente.
    """

    def __init__(self, inputs, client_logo_path, slide_width, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'client-logo', 0,
            parent_row
        )

        # Caminho da imagem no sistema
        self.client_logo_path = client_logo_path

        # Dimensões da imagem
        self.picture_width = Cm(9.53)
        self.picture_height = Cm(4.23)

    def render(self, slide):
        """
        Criar imagem no slide `pptx`.
        """

        slide.shapes.add_picture(
            self.client_logo_path,
            self.width / 2 - self.picture_width / 2,
            self.parent_row.height / 2 - self.picture_height / 2,
            self.picture_width, self.picture_height
        )


class TitleCell(Cell):
    """
    Célula responsável pelo título.
    """

    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'title', 1,
            parent_row
        )

        self.props = props

    def render(self, slide):
        """
        Renderizar.

        * slide (Slide)
        """

        # Criar retângulo para servir de fundo ao título
        shape = self.create_rect(
            Pt(-1), self.parent_row.y_offset,
            self.width + Pt(2), self.parent_row.height,
            RGBColor(0xB, 0x5D, 0x77)
        )

        # Definir a transparência do fundo como 53%
        self.set_shape_transparency(shape, 53)

        # Definir o texto do retângulo como o modelo, substituindo os espaços reservados
        # pela primeira e segunda série e a data, respetivamente
        primeira_serie = self.inputs.get('primeira-serie')
        self.set_text(
            shape,
            TITLE_STRING.format(
                primeira_serie,
                primeira_serie + 1,
                time.strftime(
                    '%B de %Y',
                    time.strptime(self.props['date'], '%d/%m/%Y')
                )
            ),
            margin_left=Cm(1)
        )


class LogoCell(Cell):
    """
    Célula responsável pelo logotipo da Logos.
    """

    def __init__(self, inputs, slide_width, parent_row):
        picture_width = Cm(6.32)
        margin = Cm(.5)

        super().__init__(
            inputs,
            {
                'width': Cm(7),
                'x_offset': slide_width - picture_width - margin
            },
            'logo', 2,
            parent_row
        )

        self.picture_width = picture_width
        self.picture_height = Cm(3.2)

    def render(self, slide):
        """
        Renderizar.

        * slide (Slide)
        """

        # Importar logo do diretório `assets` (recursos)
        with importlib.resources.path(profero.assets, 'logo_white.png') as p:
            logo_path = str(p)

        # Criar imagem com logo
        slide.shapes.add_picture(
            logo_path,
            self.x_offset + self.width / 2 - self.picture_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - self.picture_height / 2,
            self.picture_width, self.picture_height
        )


class Slide(FSlide):
    """
    Slide `title`.
    """

    def __init__(self, inputs, index, props, _, parent_presentation):
        # Importar fundo do diretório `assets`
        with importlib.resources.path(profero.assets, 'background.png') as p:
            background_path = str(p)

        super().__init__(
            inputs,
            'title', 6,
            index,
            background_path,
            parent_presentation
        )

        # Dimensões do slide
        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        # Linha do logo do cliente
        client_logo_row = Row(
            inputs,
            {
                'height': .6 * slide_height,
                'y_offset': Cm(0)
            },
            'client-logo', 0,
            self
        )

        # Célula do logo do cliente
        client_logo_cell = ClientLogoCell(
            inputs,
            inputs.get('project-logo'),
            slide_width,
            client_logo_row
        )
        client_logo_row.add_cell(client_logo_cell)

        self.add_row(client_logo_row)

        # Linha do título
        title_row = Row(
            inputs,
            {
                'height': .2 * slide_height,
                'y_offset': client_logo_row.height
            },
            'title', 1,
            self
        )

        # Célula do título
        title_cell = TitleCell(inputs, slide_width, props, title_row)
        title_row.add_cell(title_cell)

        self.add_row(title_row)

        # Linha do logo da Logos
        logo_row = Row(
            inputs,
            {
                'height': .2 * slide_height,
                'y_offset': client_logo_row.height + title_row.height
            },
            'logo', 2,
            self
        )

        # Célula do logo da Logos
        logo_cell = LogoCell(inputs, slide_width, logo_row)
        logo_row.add_cell(logo_cell)

        self.add_row(logo_row)
