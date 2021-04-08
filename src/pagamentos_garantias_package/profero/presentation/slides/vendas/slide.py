import io

from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
import matplotlib.path
import matplotlib.patches as patches

from profero.framework.presentation.slide import Slide as FSlide
from profero.framework.presentation.row import Row
from profero.framework.presentation.cell import Cell
from profero.presentation.slides.common.header import HeaderRow

import locale
locale.setlocale(locale.LC_ALL, 'pt_BR')


class ChartCell(Cell):
    def __init__(self, inputs, slide_width, x_offset, props, parent_row, index, data):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': x_offset
            },
            f'chart-{index}', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

        self.data = data

    def render(self, slide):
        month_ticks = []

        chart_width = self.width / Inches(1)
        chart_height = self.parent_row.height / Inches(1)

        plot_width = .9 if len(self.data) > 1 else .7
        plot_x = 1/2 - plot_width/2

        fig = plt.figure(figsize=(chart_width, chart_height))
        ax = fig.add_axes([plot_x, .05, plot_width, .9])

        ax.set_xlim(0, len(self.data) + 1)

        bar_width = .45 if len(self.data) > 1 else .8

        bar_displacement = 1

        colors = [
            '#895AA3',
            '#8CC257',
            '#D04042',
            '#477DBF',
        ][::-1]

        font_colors = [
            '#003960',
            '#000000',
        ]

        months = []
        max_csum = 0

        for month in self.data:
            months.append(month[0])

            values = np.array(month[:1:-1])
            csum = 0

            for (i, value), csum in zip(enumerate(values), values.cumsum()):
                if value == 0:
                    continue

                bottom = csum - value

                ax.bar(
                    bar_displacement, value,
                    bar_width, bottom=bottom,
                    color=colors[i]
                )
                plt.text(
                    bar_displacement, bottom + value / 2,
                    value,
                    ha='center',
                    va='center',
                    color=font_colors[0],
                    fontsize=9,
                    fontname='Calibri'
                )

            if csum > 0:
                plt.text(
                    bar_displacement, csum + (.7 if len(self.data) > 1 else csum * .023),
                    sum(values),
                    ha='center',
                    va='center',
                    color=font_colors[1],
                    fontsize=9,
                    fontname='Calibri'
                )

            if csum > max_csum:
                max_csum = csum

            bar_displacement += 1

        ax.ticklabel_format(useOffset=False, style='plain')
        ax.yaxis.set_major_formatter(
            FuncFormatter(lambda x, _: '{:d}'.format(x) if x != 0 else '-')
        )

        ax.tick_params(
            axis='both',
            colors='#0F3B5E',
            labelsize=7
        )

        ax.grid(
            axis='y',
            color='#eee'
        )
        ax.set_axisbelow(True)

        ax.spines['bottom'].set_color('#ddd')
        ax.spines['top'].set_color('#ddd')
        ax.spines['left'].set_color('#666')
        ax.spines['right'].set_color('#fff')

        ax.set_yticks(np.arange(0, max_csum + 1, 5 if len(self.data) > 1 else 10))
        ax.set_xticks(np.arange(1, len(months) + 1))
        ax.set_xticklabels(months)

        image_stream = io.BytesIO()
        fig.savefig(image_stream, format='png', dpi=300)

        chart_width = Inches(chart_width)
        chart_height = Inches(chart_height)

        slide.shapes.add_picture(
            image_stream,
            self.x_offset + self.width / 2 - chart_width / 2,
            self.parent_row.y_offset + self.parent_row.height / 2 - chart_height / 2,
            chart_width,
            chart_height
        )


class TableCell(Cell):
    def __init__(self, inputs, slide_width, props, parent_row):
        super().__init__(
            inputs,
            {
                'width': slide_width,
                'x_offset': 0
            },
            'table', 0,
            parent_row
        )

        self.slide_width = slide_width
        self.props = props

    def render(self, slide):
        table_width = self.slide_width * .94
        table_height = self.parent_row.height - Cm(2)
        y_correction = Cm(-.5)

        self.create_table(
            slide,
            len(self.props['sales'][0]),
            len(self.props['sales']) + 1,
            int(table_width), int(table_height),
            y_correction
        )

        self.add_table_row(
            ['Tipo'] + [month[0] for month in self.props['sales']],
            bold=True,
            color=RGBColor(255, 255, 255),
            fill_color=RGBColor(0x16, 0x36, 0x5C),
            font_size=Pt(9)
        )

        table = np.array(self.props['sales']).T
        self.add_table_row(
            table[1].astype(float).tolist(),
            'Valor',
            sub_hyphen=True,
            format_currency=True
        )

        for i, row in enumerate(table[2:]):
            self.add_table_row(
                row.astype(int),
                self.props['empreendimentos'][i],
                sub_hyphen=True,
            )


class Slide(FSlide):
    def __init__(self, inputs, index, props, table_of_contents_slide, parent_presentation):
        super().__init__(
            inputs,
            'vendas', 6,
            index,
            None,
            parent_presentation,
            'Vendas',
            table_of_contents_slide
        )

        self.props = props

        slide_height = parent_presentation.presentation.slide_height
        slide_width = parent_presentation.presentation.slide_width

        header_row = HeaderRow(
            inputs,
            {
                'height': .25 * slide_height,
                'y_offset': Cm(0)
            }, 0,
            self.title,
            slide_width, slide_height,
            self
        )
        self.add_row(header_row)

        charts_row_height = .5 * slide_height
        table_row_height = .75 * slide_height  - charts_row_height

        charts_row = Row(
            inputs,
            {
                'height': charts_row_height,
                'y_offset': header_row.y_offset + header_row.height
            },
            'charts', 1,
            self
        )

        main_chart_width = .8 * slide_width
        total_chart_width = slide_width - main_chart_width

        main_chart = ChartCell(
            inputs,
            main_chart_width, 0,
            self.props,
            charts_row, 0,
            self.props['sales']
        )
        charts_row.add_cell(main_chart)

        totals = ['Total', None]

        for month in self.props['sales']:
            for i, value in enumerate(month[2:]):
                if len(totals) < len(month):
                    totals.append(value)
                else:
                    totals[i + 2] += value

        total_chart = ChartCell(
            inputs,
            total_chart_width,
            main_chart_width,
            self.props,
            charts_row, 1,
            [totals]
        )
        charts_row.add_cell(total_chart)

        self.add_row(charts_row)

        table_row = Row(
            inputs,
            {
                'height': table_row_height,
                'y_offset': charts_row.y_offset + charts_row.height
            },
            'table', 2,
            self
        )

        table_cell = TableCell(
            self.inputs,
            slide_width,
            self.props,
            table_row
        )
        table_row.add_cell(table_cell)

        self.add_row(table_row)
